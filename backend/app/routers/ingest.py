from __future__ import annotations

from typing import List

from fastapi import APIRouter, Depends, File, Form, HTTPException, UploadFile
from sqlalchemy.orm import Session

from app.database import SessionLocal, get_db
from app.job_queue import create_job, get_job, list_jobs as list_queue_jobs, submit_local_job, update_job
from app.models.note import Note
from app.services.connection_engine import auto_connect_note
from app.services.embedding_service import add_note_embedding
from app.utils.security import validate_upload_filename

router = APIRouter(prefix="/api/ingest", tags=["ingest"])


def _new_job(filename: str) -> str:
    return create_job("ingest-file", {"filename": filename})["id"]


def _update_job(job_id: str, **kwargs):
    update_job(job_id, **kwargs)


def _extract_text(filename: str, data: bytes) -> str:
    lower = filename.lower()
    if lower.endswith((".txt", ".md")):
        return data.decode("utf-8", errors="replace")
    if lower.endswith((".html", ".htm")):
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(data, "html.parser")
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        return soup.get_text(separator=" ", strip=True)
    if lower.endswith(".pdf"):
        try:
            import io
            import pdfplumber

            with pdfplumber.open(io.BytesIO(data)) as pdf:
                return "\n\n".join(page.extract_text() or "" for page in pdf.pages).strip()
        except ImportError as exc:
            raise HTTPException(status_code=422, detail="PDF parsing requires pdfplumber") from exc
    if lower.endswith(".docx"):
        try:
            import io
            import docx

            doc = docx.Document(io.BytesIO(data))
            return "\n".join(item.text for item in doc.paragraphs)
        except ImportError as exc:
            raise HTTPException(status_code=422, detail="DOCX parsing requires python-docx") from exc
    if lower.endswith(".pptx"):
        try:
            import io
            from pptx import Presentation

            presentation = Presentation(io.BytesIO(data))
            lines: List[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        lines.append(shape.text)
            return "\n".join(lines)
        except ImportError as exc:
            raise HTTPException(status_code=422, detail="PPTX parsing requires python-pptx") from exc
    raise HTTPException(status_code=415, detail="Unsupported file type")


def _process_file_job(job_id: str, filename: str, raw_bytes: bytes, note_title: str, db_session: Session):
    _update_job(job_id, status="extracting", progress=20)
    try:
        text = _extract_text(filename, raw_bytes)
        if not text.strip():
            raise HTTPException(status_code=422, detail="No text extracted from file")

        _update_job(job_id, status="saving", progress=45)
        note = Note(
            title=(note_title or filename)[:200],
            content=f"**Source:** {filename}\n\n{text[:12000]}",
            source=filename,
            source_type="file",
            tags="[]",
        )
        db_session.add(note)
        db_session.commit()
        db_session.refresh(note)
        _update_job(job_id, note_id=note.id, status="embedding", progress=70)

        try:
            add_note_embedding(note.id, note.title, note.content)
        except Exception as exc:
            _update_job(job_id, embedding_error=str(exc))

        _update_job(job_id, status="connecting", progress=85)
        try:
            auto_connect_note(note.id)
        except Exception as exc:
            _update_job(job_id, connection_error=str(exc))

        _update_job(job_id, status="done", progress=100, result={"note_id": note.id})
    except HTTPException as exc:
        _update_job(job_id, status="failed", error=exc.detail)
    except Exception as exc:
        _update_job(job_id, status="failed", error=str(exc))
    finally:
        db_session.close()


@router.post("/file")
async def upload_file(file: UploadFile = File(...), title: str = Form(""), db: Session = Depends(get_db)):
    del db
    data = await file.read()
    if len(data) > 25 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="File too large (max 25 MB)")
    filename = file.filename or "upload"
    if not validate_upload_filename(filename):
        raise HTTPException(status_code=415, detail="Unsupported or unsafe filename extension")

    job_id = _new_job(filename)
    submit_local_job(job_id, _process_file_job, job_id, filename, data, title.strip(), SessionLocal())
    return {"job_id": job_id, "filename": filename, "status": "queued"}


@router.get("/jobs")
def list_jobs():
    return list_queue_jobs(50)


@router.get("/jobs/{job_id}")
def get_job_status(job_id: str):
    job = get_job(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found")
    return job


@router.delete("/jobs/{job_id}")
def dismiss_job(job_id: str):
    update_job(job_id, dismissed=True)
    return {"ok": True}
